#!/usr/bin/env python3
"""docgen.py <kind> <outfile> -- lee de stdin un texto en markdown liviano y lo convierte
a un documento real con formato (docx/pdf/pptx/xlsx).

Convenciones del markdown liviano (las mismas para los 4 formatos):
  # Titulo            -> encabezado nivel 1
  ## Subtitulo        -> encabezado nivel 2
  - item               -> vinieta
  linea suelta         -> parrafo normal
  ---                  -> separador de slide (SOLO para pptx: separa una slide de la
                          siguiente; la primera '#'/'##' de cada bloque es el titulo
                          de esa slide y las '- ' son sus vinietas)
  a,b,c (o a|b|c)      -> fila de tabla (SOLO para xlsx: cada linea es una fila, separada
                          por coma o por pipe; se usa la primera fila como encabezado)
  !img <que buscar>    -> busca una imagen LIBRE en Wikimedia Commons, la baja y la inserta ahi
                          mismo, con la atribucion (autor y licencia) como epigrafe.
                          Ej:  !img panel solar en un techo
  !imgfile <ruta>      -> inserta una imagen que ya esta en el disco (por ejemplo una recien
                          generada). Se le puede poner epigrafe con '|':
                          Ej:  !imgfile C:/ruta/foto.jpg|Ilustracion generada
"""
import sys


# --- IMAGENES EN LOS DOCUMENTOS QUE MENTIS GENERA (2026-08-03, B2 del plan) ----------------------
#
# Sintaxis:  !img <que buscar>
#
# La imagen se busca en Wikimedia Commons y se baja al disco ANTES de armar el documento (ver
# _resolver_imagenes mas abajo). Dos decisiones que no son de estilo:
#
#   1. SE BAJA, NO SE ENLAZA. Un documento con imagenes enlazadas se rompe solo: el dia que el
#      sitio cambia la URL queda un recuadro roto, meses despues, cuando ya nadie se acuerda.
#   2. SIEMPRE VA LA ATRIBUCION, como epigrafe debajo de la imagen. Todo lo de Commons es de
#      licencia libre, pero libre no es lo mismo que sin autor. Poner la foto de alguien en un
#      documento del usuario sin decir de quien es le crearia un problema en vez de resolverle uno.
#
# Si la busqueda no encuentra nada, NO se falla el documento entero ni se borra la linea en
# silencio: queda un parrafo diciendo que ahi iba una imagen y cual era. Un hueco silencioso es
# peor que un hueco senalado -- nadie revisa lo que no sabe que falta.
def parse_blocks(text):
    lines = text.splitlines()
    blocks = []
    for ln in lines:
        s = ln.rstrip()
        if not s.strip():
            blocks.append(("blank", ""))
        elif s.startswith("!imgok "):
            # Ya resuelta por _resolver_imagenes: "ruta|atribucion"
            resto = s[7:].strip()
            ruta, _, atrib = resto.partition("|")
            blocks.append(("imgok", (ruta.strip(), atrib.strip())))
        elif s.startswith("!imgfile "):
            # Una imagen que YA esta en el disco (por ejemplo una recien generada con gen).
            #
            # ESTE CASO FALTABA Y SE NOTO PROBANDO. Mentis generaba una imagen y despues escribia
            # en el documento un texto '[IMAGEN: archivo.jpg]', que no inserta nada: un cartel
            # avisando que el documento quedo sin la imagen. No era terquedad del modelo -- era
            # que no habia ninguna sintaxis para insertar un archivo local, asi que el cartel era
            # su unica opcion.
            resto = s[9:].strip()
            ruta, _, atrib = resto.partition("|")
            blocks.append(("imgok", (ruta.strip(), atrib.strip())))
        elif s.startswith("!img "):
            # Llego sin resolver (por ejemplo si alguien llama a docgen.py directo): no se
            # inventa nada ni se borra la linea; queda dicho que ahi faltaba una imagen.
            blocks.append(("p", "[falta la imagen: %s]" % s[5:].strip()))
        elif s.startswith("## "):
            blocks.append(("h2", s[3:].strip()))
        elif s.startswith("# "):
            blocks.append(("h1", s[2:].strip()))
        elif s.startswith("- "):
            blocks.append(("bullet", s[2:].strip()))
        else:
            blocks.append(("p", s.strip()))
    return blocks


def _resolver_imagenes(text, destdir):
    """Cambia cada '!img <consulta>' por '!imgok <ruta>|<atribucion>' bajando la imagen.

    Se hace en una pasada previa y no dentro de cada generador para no repetir la descarga cuatro
    veces, y para que los cuatro formatos se comporten igual.
    """
    import os
    import subprocess
    aqui = os.path.dirname(os.path.abspath(__file__))
    buscador = os.path.join(aqui, "engine", "img_buscar.py")
    if not os.path.isfile(buscador):
        buscador = os.path.join(aqui, "img_buscar.py")
    salida = []
    n = 0
    for linea in text.splitlines():
        if not linea.startswith("!img "):
            salida.append(linea)
            continue
        consulta = linea[5:].strip()
        n += 1
        if not consulta or not os.path.isfile(buscador):
            salida.append("[falta la imagen: %s]" % (consulta or "sin consulta"))
            continue
        try:
            r = subprocess.run([sys.executable, buscador, consulta, "--dest", destdir,
                                "--nombre", "doc-img-%d" % n, "--json"],
                               capture_output=True, text=True, timeout=90, encoding="utf-8")
            import json as _json
            d = _json.loads((r.stdout or "").strip().splitlines()[0])
            salida.append("!imgok %s|%s" % (d["ruta"], d["atribucion"]))
        except Exception:
            salida.append("[no se pudo conseguir una imagen de: %s]" % consulta)
    return "\n".join(salida)


def gen_docx(text, outfile):
    from docx import Document
    from docx.shared import Pt
    doc = Document()
    for kind, content in parse_blocks(text):
        if kind == "blank":
            continue
        elif kind == "h1":
            doc.add_heading(content, level=1)
        elif kind == "h2":
            doc.add_heading(content, level=2)
        elif kind == "bullet":
            doc.add_paragraph(content, style="List Bullet")
        elif kind == "imgok":
            ruta, atrib = content
            try:
                from docx.shared import Inches
                doc.add_picture(ruta, width=Inches(5.5))
                # La atribucion va SIEMPRE, como epigrafe. Commons es de licencia libre, pero
                # libre no es sin autor: la foto de alguien en un documento del usuario sin decir de
                # quien es le crea un problema en vez de resolverle uno.
                cap = doc.add_paragraph(atrib)
                cap.runs[0].italic = True
                cap.runs[0].font.size = Pt(8)
            except Exception as e:
                doc.add_paragraph("[no se pudo insertar la imagen: %s]" % (str(e)[:80] or type(e).__name__))
        else:
            doc.add_paragraph(content)
    doc.save(outfile)


def gen_pdf(text, outfile):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem
    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(outfile, pagesize=A4)
    story = []
    bullets = []

    def flush_bullets():
        nonlocal bullets
        if bullets:
            story.append(ListFlowable([ListItem(Paragraph(b, styles["Normal"])) for b in bullets], bulletType="bullet"))
            bullets = []

    for kind, content in parse_blocks(text):
        if kind == "blank":
            flush_bullets()
            story.append(Spacer(1, 8))
        elif kind == "h1":
            flush_bullets()
            story.append(Paragraph(content, styles["Heading1"]))
        elif kind == "h2":
            flush_bullets()
            story.append(Paragraph(content, styles["Heading2"]))
        elif kind == "bullet":
            bullets.append(content)
        elif kind == "imgok":
            flush_bullets()
            ruta, atrib = content
            try:
                from reportlab.platypus import Image as RLImage
                from reportlab.lib.utils import ImageReader
                iw, ih = ImageReader(ruta).getSize()
                ancho = min(400.0, float(iw))
                story.append(RLImage(ruta, width=ancho, height=ancho * ih / float(iw)))
                story.append(Paragraph("<i><font size=7>%s</font></i>" % atrib, styles["Normal"]))
            except Exception as e:
                story.append(Paragraph("[no se pudo insertar la imagen: %s]" % (str(e)[:80] or type(e).__name__), styles["Normal"]))
        else:
            flush_bullets()
            story.append(Paragraph(content, styles["Normal"]))
    flush_bullets()
    doc.build(story)


def gen_pptx(text, outfile):
    from pptx import Presentation
    prs = Presentation()
    title_layout = prs.slide_layouts[1]  # Title and Content
    raw_slides = text.split("\n---\n") if "\n---\n" in text else [text]
    for raw in raw_slides:
        blocks = parse_blocks(raw)
        title = ""
        bullets = []
        for kind, content in blocks:
            if kind in ("h1", "h2") and not title:
                title = content
            elif kind == "bullet":
                bullets.append(content)
            elif kind == "p" and content and not title:
                title = content
            elif kind == "p" and content:
                bullets.append(content)
        imagenes = [c for k, c in blocks if k == "imgok"]
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = title or "(sin titulo)"
        body = slide.placeholders[1].text_frame
        body.clear()
        for i, b in enumerate(bullets):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = b
        # La imagen va DEBAJO del cuerpo y ocupando el ancho util, no como una vinieta mas: en una
        # presentacion la imagen es el contenido, no una nota al pie.
        for ruta, atrib in imagenes[:1]:
            try:
                from pptx.util import Inches, Pt as _Pt
                slide.shapes.add_picture(ruta, Inches(1.2), Inches(3.2), width=Inches(7.0))
                caja = slide.shapes.add_textbox(Inches(1.2), Inches(6.6), Inches(7.0), Inches(0.4))
                run = caja.text_frame.paragraphs[0].add_run()
                run.text = atrib
                run.font.size = _Pt(8)
                run.font.italic = True
            except Exception:
                pass
    prs.save(outfile)


def gen_xlsx(text, outfile):
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    for ln in text.splitlines():
        s = ln.strip()
        if not s:
            continue
        sep = "|" if "|" in s else ","
        row = [c.strip() for c in s.split(sep)]
        ws.append(row)
    wb.save(outfile)


def main():
    if len(sys.argv) != 3:
        print("Uso: docgen.py <docx|pdf|pptx|xlsx> <archivo_salida>", file=sys.stderr)
        sys.exit(1)
    kind, outfile = sys.argv[1], sys.argv[2]
    text = sys.stdin.read()
    # Las imagenes se bajan ANTES de armar el documento, una sola vez, y al lado del archivo de
    # salida: asi quedan junto al documento que las usa y no en un temporal que se limpia solo.
    if "!img " in text:
        import os
        destino = os.path.join(os.path.dirname(os.path.abspath(outfile)), "medios")
        text = _resolver_imagenes(text, destino)
    fn = {"docx": gen_docx, "pdf": gen_pdf, "pptx": gen_pptx, "xlsx": gen_xlsx}.get(kind)
    if not fn:
        print(f"ERROR: tipo desconocido '{kind}'. Usa docx|pdf|pptx|xlsx.", file=sys.stderr)
        sys.exit(1)
    fn(text, outfile)
    print(outfile)


if __name__ == "__main__":
    main()
