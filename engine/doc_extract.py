"""doc_extract.py -- saca el texto Y las imagenes de un documento (2026-08-03, B1 del plan).

EL PROBLEMA: hasta hoy la tool 'read' rechazaba de plano cualquier archivo binario. Un.docx ES
binario, asi que Mentis no podia leer un Word en absoluto -- ni su texto. Y cuando el usuario le pasaba
un informe con graficos, las imagenes no existian para ella: contestaba sobre el texto como si el
documento no tuviera nada mas.

QUE HACE: devuelve el texto en orden y, donde habia una imagen, deja una marca con su ruta:

    Parrafo de texto...
    [[IMAGEN 1: /tmp/xxx/img-1.png]]
    Mas texto...

Quien llama decide que hacer con las imagenes (nv-agent.sh se las pasa al rol multimodal y
reemplaza cada marca por la descripcion). Aca no se llama a ningun modelo: esto es extraccion
pura, deterministica y gratis.

FORMATOS:.docx,.pptx,.xlsx,.pdf. Todo con librerias que YA estaban instaladas en esta
maquina (python-docx, python-pptx, openpyxl, PyMuPDF) -- no se agrego ninguna dependencia.

Uso:  doc_extract.py <archivo> [--imgdir <carpeta>] [--max-img N] [--json]
"""
import argparse
import io
import json
import os
import sys
import zipfile

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8", errors="replace")

# Debajo de esto no vale la pena: son vinetas, lineas divisorias, logos de 1x1. Describir una
# imagen cuesta una llamada al modelo multimodal, y gastarla en un separador es tirarla.
MIN_BYTES = 3000
MAX_IMG_DEFECTO = 6


def _guardar(imgdir, n, datos, ext):
    # Barras NORMALES siempre, aunque estemos en Windows. Quien consume estas rutas es bash
    # (nv-agent.sh se las pasa a ask-nvidia.sh -I), y os.path.join en Windows devuelve
    # "C:/Users/...\img-1.png" -- mezcla que bash no sabe abrir. Con barras uniformes andan tanto
    # bash como cygpath como el propio Python.
    ruta = imgdir.replace("\\", "/").rstrip("/") + "/img-%d%s" % (n, ext)
    with open(ruta, "wb") as f:
        f.write(datos)
    return ruta


def _ext_de(nombre, defecto=".png"):
    e = os.path.splitext(nombre)[1].lower()
    return e if e in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp") else defecto


# --- docx -----------------------------------------------------------------------------------------
def extraer_docx(ruta, imgdir, tope):
    from docx import Document
    from docx.oxml.ns import qn

    doc = Document(ruta)
    partes = []
    imgs = []

    # Se recorre el XML del cuerpo en ORDEN, no los parrafos por un lado y las imagenes por otro.
    # Si se listan aparte se pierde donde estaba cada imagen, y "el grafico de abajo" deja de
    # tener sentido -- que es justo lo que el usuario quiere poder preguntar.
    rels = doc.part.rels
    for bloque in doc.element.body.iter():
        if bloque.tag == qn("w:p"):
            texto = "".join(n.text or "" for n in bloque.iter(qn("w:t")))
            if texto.strip():
                partes.append(texto)
            for blip in bloque.iter(qn("a:blip")):
                rid = blip.get(qn("r:embed"))
                if not rid or rid not in rels:
                    continue
                try:
                    blob = rels[rid].target_part.blob
                except Exception:
                    continue
                if len(blob) < MIN_BYTES or len(imgs) >= tope:
                    continue
                nombre = getattr(rels[rid].target_part, "partname", "img.png")
                p = _guardar(imgdir, len(imgs) + 1, blob, _ext_de(str(nombre)))
                imgs.append(p)
                partes.append("[[IMAGEN %d: %s]]" % (len(imgs), p))
        elif bloque.tag == qn("w:tbl"):
            filas = []
            for tr in bloque.iter(qn("w:tr")):
                celdas = ["".join(n.text or "" for n in tc.iter(qn("w:t"))).strip()
                          for tc in tr.iter(qn("w:tc"))]
                if any(celdas):
                    filas.append(" | ".join(celdas))
            if filas:
                partes.append("TABLA:\n" + "\n".join(filas))
    return "\n".join(partes), imgs


# --- pptx -----------------------------------------------------------------------------------------
def extraer_pptx(ruta, imgdir, tope):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    pres = Presentation(ruta)
    partes = []
    imgs = []
    for i, slide in enumerate(pres.slides, 1):
        partes.append("--- Diapositiva %d ---" % i)
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                partes.append(shape.text_frame.text.strip())
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and len(imgs) < tope:
                try:
                    blob = shape.image.blob
                    ext = "." + (shape.image.ext or "png")
                except Exception:
                    continue
                if len(blob) < MIN_BYTES:
                    continue
                p = _guardar(imgdir, len(imgs) + 1, blob, _ext_de("x" + ext))
                imgs.append(p)
                partes.append("[[IMAGEN %d: %s]]" % (len(imgs), p))
    return "\n".join(partes), imgs


# --- xlsx -----------------------------------------------------------------------------------------
def extraer_xlsx(ruta, imgdir, tope):
    import openpyxl
    wb = openpyxl.load_workbook(ruta, data_only=True)
    partes = []
    for hoja in wb.worksheets:
        partes.append("--- Hoja: %s ---" % hoja.title)
        for fila in hoja.iter_rows(values_only=True):
            if any(c is not None and str(c).strip() for c in fila):
                partes.append(" | ".join("" if c is None else str(c) for c in fila))
    # openpyxl no expone las imagenes de forma estable entre versiones; se sacan del zip, que es
    # lo que un.xlsx es por dentro.
    imgs = []
    try:
        with zipfile.ZipFile(ruta) as z:
            for nombre in sorted(n for n in z.namelist() if n.startswith("xl/media/")):
                if len(imgs) >= tope:
                    break
                datos = z.read(nombre)
                if len(datos) < MIN_BYTES:
                    continue
                p = _guardar(imgdir, len(imgs) + 1, datos, _ext_de(nombre))
                imgs.append(p)
                partes.append("[[IMAGEN %d: %s]]" % (len(imgs), p))
    except Exception:
        pass
    return "\n".join(partes), imgs


# --- pdf ------------------------------------------------------------------------------------------
def extraer_pdf(ruta, imgdir, tope):
    import fitz
    doc = fitz.open(ruta)
    partes = []
    imgs = []
    vistos = set()
    for n, pagina in enumerate(doc, 1):
        partes.append("--- Pagina %d ---" % n)
        t = pagina.get_text().strip()
        if t:
            partes.append(t)
        for info in pagina.get_images(full=True):
            if len(imgs) >= tope:
                break
            xref = info[0]
            if xref in vistos:      # la misma imagen repetida en varias paginas (logos, marcas)
                continue
            vistos.add(xref)
            try:
                base = doc.extract_image(xref)
            except Exception:
                continue
            datos = base.get("image") or b""
            if len(datos) < MIN_BYTES:
                continue
            p = _guardar(imgdir, len(imgs) + 1, datos, "." + (base.get("ext") or "png"))
            imgs.append(p)
            partes.append("[[IMAGEN %d: %s]]" % (len(imgs), p))
    doc.close()
    return "\n".join(partes), imgs


EXTRACTORES = {".docx": extraer_docx, ".pptx": extraer_pptx,
               ".xlsx": extraer_xlsx, ".pdf": extraer_pdf}


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("archivo")
    ap.add_argument("--imgdir", default=None)
    ap.add_argument("--max-img", type=int, default=MAX_IMG_DEFECTO)
    ap.add_argument("--json", action="store_true")
    args = ap.parse_args()

    if not os.path.isfile(args.archivo):
        sys.stderr.write("no existe: %s\n" % args.archivo)
        return 1
    ext = os.path.splitext(args.archivo)[1].lower()
    fn = EXTRACTORES.get(ext)
    if not fn:
        sys.stderr.write("formato no soportado: %s (van.docx.pptx.xlsx.pdf)\n" % ext)
        return 2

    imgdir = args.imgdir or os.path.join(os.path.dirname(os.path.abspath(args.archivo)),
                                         ".mentis-medios")
    os.makedirs(imgdir, exist_ok=True)

    try:
        texto, imgs = fn(args.archivo, imgdir, max(0, args.max_img))
    except Exception as e:
        sys.stderr.write("no pude leer el documento: %s: %s\n" % (type(e).__name__, str(e)[:200]))
        return 3

    if args.json:
        print(json.dumps({"texto": texto, "imagenes": imgs, "formato": ext},
                         ensure_ascii=False))
    else:
        print(texto)
        if imgs:
            sys.stderr.write("IMAGENES=%s\n" % "|".join(imgs))
    return 0


if __name__ == "__main__":
    sys.exit(main())
